"""
Document processing pipeline (docs/ARCHITECTURE.md §4 step-by-step).
Called from a Celery task, never from the API request thread directly.
"""
import uuid
import tempfile
import os
import boto3
from botocore.client import Config

from app.core.config import get_settings
from app.core.llm_gateway import get_provider

settings = get_settings()

s3_client = boto3.client(
    "s3",
    endpoint_url=settings.S3_ENDPOINT_URL,
    aws_access_key_id=settings.S3_ACCESS_KEY,
    aws_secret_access_key=settings.S3_SECRET_KEY,
    config=Config(signature_version="s3v4"),
)


async def extract_text(storage_path: str, file_type: str) -> str:
    """
    Step 2: extract raw text from the stored file in S3/MinIO.
    Supports PDF, DOCX, PPTX, TXT, MD.
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
        temp_path = tmp.name

    try:
        # Download from S3/MinIO
        s3_client.download_file(settings.S3_BUCKET, storage_path, temp_path)

        text = ""
        if file_type == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(temp_path)
            parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    parts.append(page_text)
            text = "\n".join(parts)

        elif file_type == "docx":
            import docx
            doc = docx.Document(temp_path)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(temp_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            text = "\n".join(parts)

        elif file_type in ("txt", "md"):
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        else:
            raise NotImplementedError(f"Extraction not yet implemented for file type: {file_type}")

        return text
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


def clean_text(raw_text: str) -> str:
    """Step 3: normalize whitespace, strip boilerplate/headers-footers, fix encoding artifacts."""
    return " ".join(raw_text.split())


def chunk_text(text: str, strategy: str = "fixed", chunk_size: int = 800, overlap: int = 120) -> list[dict]:
    """
    Step 4: split into chunks.
      - fixed: sliding window over tokens/characters
      - semantic: split on topic/sentence-embedding similarity shifts
      - parent_child: small child chunks for retrieval, larger parent chunk for context injection
    Returns list of {content, page_number, paragraph_index}.
    """
    if strategy == "fixed":
        words = text.split()
        chunks = []
        step = chunk_size - overlap
        if step <= 0:
            step = chunk_size
        for i in range(0, len(words), step):
            piece = " ".join(words[i : i + chunk_size])
            if piece:
                chunks.append({"content": piece, "page_number": None, "paragraph_index": i // step})
        return chunks
    raise NotImplementedError(f"Chunk strategy '{strategy}' not yet implemented")


async def embed_chunks(chunks: list[dict], provider_name: str | None = None) -> list[dict]:
    """Step 5: generate embeddings for each chunk's content."""
    provider = get_provider(provider_name)
    texts = [c["content"] for c in chunks]
    vectors = await provider.embed(texts)
    for chunk, vector in zip(chunks, vectors):
        chunk["embedding"] = vector
    return chunks


async def process_document(document_id: uuid.UUID) -> None:
    """
    Orchestrates steps 2-8 for a single document.
    Downloads the file, extracts text, cleans it, chunks it, embeds it, and indexes the vectors.
    """
    from app.db.session import AsyncSessionLocal, engine
    from app.models.documents import Document, DocumentStatus, Chunk
    from sqlalchemy import select

    # Dispose of existing engine connection pool to avoid event loop mismatch conflicts in workers
    await engine.dispose()

    async with AsyncSessionLocal() as db:
        # 1. Load Document row
        stmt = select(Document).where(Document.id == document_id)
        res = await db.execute(stmt)
        doc = res.scalar_one_or_none()
        if not doc:
            raise ValueError(f"Document with ID {document_id} not found in database")

        # 2. Mark status=processing
        doc.status = DocumentStatus.processing
        await db.commit()

        try:
            # 3. Extract text
            raw_text = await extract_text(doc.storage_path, doc.file_type)

            # 4. Clean text
            cleaned = clean_text(raw_text)
            if not cleaned:
                raise ValueError("Extracted text is empty or invalid.")

            # 5. Chunk text
            chunks = chunk_text(cleaned)

            # 6. Embed chunks
            chunks_with_embeddings = await embed_chunks(chunks)

            # 7. Bulk insert Chunk rows
            db_chunks = [
                Chunk(
                    id=uuid.uuid4(),
                    document_id=doc.id,
                    content=c["content"],
                    embedding=c["embedding"],
                    page_number=c.get("page_number"),
                    paragraph_index=c.get("paragraph_index"),
                )
                for c in chunks_with_embeddings
            ]
            db.add_all(db_chunks)

            # 8. Mark status=ready
            doc.status = DocumentStatus.ready
            await db.commit()

        except Exception as e:
            # Mark status=failed on exception
            doc.status = DocumentStatus.failed
            await db.commit()
            raise e
